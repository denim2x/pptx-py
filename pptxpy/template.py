# encoding: utf-8

from posixpath import splitext
from pptx import Presentation as _load
from .common import lazyproperty, Slides, RT, Cache, Presentation, PresentationPart, PackURI, Slide, Part, Rel
from .slide import _Slide


class Template:
  def __init__(self, uri):
    self._uri = uri
    self._model = None

  @property
  def model(self):
    return self._model

  def __call__(self):
    prs = _load(self._uri)     # FIXME: Check _uri existence
    if self._model is None:
      self._model = _Model(prs.slides)

    for node in prs.slide_layouts:
      node.part.drop_all(RT.SLIDE, recursive=False)

    prs.part.drop_all(RT.SLIDE)
    prs.slides.clear()    

    return _Presentation(prs.part, self)

  def __len__(self):
    if self._model is None:
      return

    return len(self._model)


class _Presentation(Presentation):
  def __init__(self, part, template):
    Presentation.__init__(self, part._element, part)
    part._presentation = self
    self._source = template

  @property
  def model(self):
    return self._source.model

  def save(self, file, update_links=False):  # FIXME: Clean extraneous relationships (including those without corresponding links)
    for slide in self.slides:
      slide._update(self.slides, update_links)

    self.part.save(file)
    return self

  @lazyproperty
  def slides(self):
    sldIdLst = self._element.get_or_add_sldIdLst()
    self.part.rename_slide_parts([sldId.rId for sldId in sldIdLst])
    return _Slides(sldIdLst, self)  


class _Slides(Slides):
  def __init__(self, sldIdLst, prs):
    Slides.__init__(self, sldIdLst, prs)
    self._model = prs.model

  def __call__(self, slide_index=None, slide_id=None, position=None):
    model = self._model
    if model is None:
      return

    if slide_index is None and slide_id is None:
      ret = []

      for i in range(len(model)):
        ret.append(self(i))

      return ret

    slide_model = model(slide_index, slide_id)
    if slide_model is None:
      return

    ret = slide_model(Cache(self.part.package))
    rId = self.part.relate_to(ret, RT.SLIDE)

    slide_id = self._sldIdLst.add_sldId(rId)
    slide = _Slide(ret, slide_id)
    slide.relocate(position)
    return slide



class _Model:
  def __init__(self, slides):
    self._list = _SlideParts(slides, self)
    self._ids = {}

  def get(self, slide_id):
    if slide_id not in self._ids:
      for item in self._list:
        if item.slide_id == slide_id:
          self._ids[slide_id] = item
          break

    return self._ids[slide_id]

  def __getitem__(self, slide_index):
    return self._list[slide_index]

  def __iter__(self):
    return self._list.values()

  def __len__(self):
    return len(self._list)

  def __call__(self, slide_index=None, slide_id=None):
    """
    Retrieve the |_Slide| instance from the given by either
    (but not both) *slide_index* or *slide_id*.
    """
    model = None
    if slide_index is not None:
      model = self[slide_index]
    elif slide_id is not None:
      model = self.get(slide_id)

    return model


class _SlideParts(object):
  def __new__(cls, source, owner=None):
    return [_Part(s.part, owner, s.slide_id) for s in source]


class _Reference:
  def __init__(self, part, owner=None, slide_id=None):
    base, ext = splitext(part.partname)
    self._partname = PackURI('%s%s' % (base, ext.lower()) if ext else base)
    self._content_type = part.content_type
    self._blob = part.blob
    self._owner = owner
    self._slide_id = slide_id

  @property
  def blob(self):
    return self._blob

  @property
  def slide_id(self):
    return self._slide_id

  @property
  def partname(self):
    return self._partname

  @property
  def content_type(self):
    return self._content_type


class _Part(_Reference):
  def __init__(self, part, owner=None, slide_id=None):
    _Reference.__init__(self, part, owner, slide_id)
    part._model = self
    self._uri = self.partname.template
    self._load = part.load    
    self._rels = _Relationships(part, self)

  def __call__(self, cache):
    uri = cache.next_partname(self._uri)
    part = self._load(uri, self._content_type, self._blob, cache.package)

    cache[self] = part

    for rel in self._rels:
      if rel.reltype == RT.SLIDE:
        continue

      out = rel(part, cache)
      if rel.reltype in part._reltypes:
        out.target_part.relate_to(part, part._reltype)

    return part

  @classmethod
  def get(cls, part, owner):
    if hasattr(part, '_model'):
      return part._model

    return cls(part, owner)


class _Relationships(object):
  def __new__(cls, part, owner=None):
    return [_Relationship(rel, owner) for rel in part.rels.values()]


class _Relationship:
  def __init__(self, rel, owner=None):
    self._reltype = rel.reltype
    self._is_external = rel.is_external
    self._rId = rel.rId
    if self.reltype == RT.SLIDE:
      self._target = _Reference(rel.target_part, owner)
    elif not self.is_external:
      self._target = _Part.get(rel.target_part, owner)
    else:
      self._target = rel.target_ref

  @property
  def is_external(self):
    return self._is_external

  @property
  def reltype(self):
    return self._reltype

  @property
  def target(self):
    return self._target  

  def __call__(self, part, cache=None):
    target = self.target
    if not self.is_external and self.reltype != RT.SLIDE:
      target = part.package[self.target, self.reltype]

      if target is None and cache is not None:
        target = cache[self.target]

    return part.load_rel(self.reltype, target, self._rId, self.is_external)
