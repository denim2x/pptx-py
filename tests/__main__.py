import os
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

dirname = os.path.dirname(__file__)
try:
  import pptxy
except ImportError:
  import sys
  sys.path.insert(0, os.path.normpath(os.path.join(dirname, '..')))
  import pptxpy


def normpath(path):
  return os.path.normpath(os.path.join(dirname, path))

prs = Presentation(normpath('test_files/test_slides.pptx'))
slide_master1 = prs.slide_masters[0]


def test_duplicate():
  global prs
  num = len(prs.slides)
  assert num > 0, "Cannot test Presentation with no slides"

  for i in range(num):
    _duplicate(i)


def test_remove():
  global prs
  num = len(prs.slides)
  assert num > 0, "Cannot test Presentation with no slides"

  s = prs.slides.remove(0)

  assert len(prs.slides) == num - 1
  assert len(prs.slides) == 0 or prs.slides[0] is not s

test_remove()

import sys
if len(sys.argv) > 1:
  path = sys.argv[1]
  prs.save(path)


def _duplicate(i):
  global prs, slide_master1
  num_rels = len(prs.part.rels)

  s, num = prs.slides[i], len(prs.slides)
  l = s.slide_layout.part
  m = l.slide_master

  c = prs.slides.duplicate(i, slide_master=m is slide_master1)
  assert c.slide_layout.part is not l, "Slide #%d's SlideLayout wasn't cloned" % i

  if m is slide_master1:
    assert c.slide_layout.part.slide_master is not m, "Slide #%d's SlideMaster wasn't cloned" % i
  else:
    assert c.slide_layout.part.slide_master is m, "Slide #%d's SlideMaster is not OK" % i

  assert len(prs.slides) == num + 1
  assert prs.slides[-1] is c

  sp, cp = s.part, c.part
  assert sp.partname.is_similar(cp.partname)
  assert sp.content_type == cp.content_type
  assert sp.blob == cp.blob
  assert sp.package == cp.package
  assert len(prs.part.rels) > num_rels    # FIXME

  return
  assert sp.rels.equals(cp.rels, False)   # FIXME

  assert sp.rels == cp.rels, \
    'slides[%d].rels != slides[%d].rels (%s != %s)' % (
      i, num, sp.rels.pprint(), cp.rels.pprint()
    )


def test_background():
  a_solidFills = prs.slides[0].background.element.xpath('//a:solidFill')
  c_solidFills = c.background.element.xpath('//a:solidFill')

  c.background.fill.solid()

  assert prs.slides[0].part.blob != c.part.blob
  assert c.background.fill.type == 1
